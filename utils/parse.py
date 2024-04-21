from pptx import Presentation
from PyPDF2 import PdfReader

def parse_text(file):
    extracted_text = ""
    # Checks if file is a powerpoint
    if file.filename.endswith('.pptx'):
        # Converts file to Powerpoint object
        prs = Presentation(file)
        # Loops through slides
        for slide_number, slide in enumerate(prs.slides):
            # Adds slide number before the text in each slide
            extracted_text += f"\nSlide {slide_number + 1}:\n"
            # Parses text in each slide + adds it
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
    # Checks if file is a pdf
    elif file.filename.endswith('.pdf'):
        # Creates pdf reader for file
        reader = PdfReader(file)
        # Loops through pages in pdf, extracting text
        for page in reader.pages:
            extracted_text += page.extract_text() 

    return extracted_text

