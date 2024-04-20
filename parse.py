from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    extracted_text = ""
    for slide_number, slide in enumerate(presentation.slides):
        extracted_text += f"\nSlide {slide_number + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
    return extracted_text

# Replace 'your_presentation.pptx' with the path to your PowerPoint file
pptx_path = 'Grocer.pptx'
text = extract_text_from_pptx(pptx_path)

# Display or save the extracted text as needed
print(text)
