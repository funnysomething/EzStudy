from pptx import Presentation
from PyPDF2 import PdfReader

def parse_text(file):
    extracted_text = ""
    if file.filename.endswith('.pptx'):
        prs = Presentation(file)
        for slide_number, slide in enumerate(prs.slides):
            extracted_text += f"\nSlide {slide_number + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
    elif file.filename.endswith('.pdf'):
        reader = PdfReader(file)
        for page in reader.pages:
            extracted_text += page.extract_text() 

    return extracted_text

